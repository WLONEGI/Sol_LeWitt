import json
import logging
import os
import subprocess
import tempfile
import zipfile
from pathlib import Path
from typing import Annotated

from langchain_core.tools import tool
from pdf2image import convert_from_path
from PIL import Image
from pydantic import BaseModel, Field

from src.domain.designer.pptx_parser import extract_pptx_context
from src.domain.designer.pdf import assemble_pdf_from_images

logger = logging.getLogger(__name__)


def _resolve_path(path: str, work_dir: str) -> str:
    candidate = Path(path)
    if candidate.is_absolute():
        return str(candidate)
    return str(Path(work_dir) / candidate)


def _ensure_dir(path: str) -> str:
    os.makedirs(path, exist_ok=True)
    return path


def _build_pptx_from_images(image_paths: list[str], output_path: str, title: str | None = None) -> None:
    try:
        from pptx import Presentation
        from pptx.util import Emu
    except Exception as exc:
        raise RuntimeError(f"python-pptx is required for PPTX packaging: {exc}") from exc

    if not image_paths:
        raise ValueError("image_paths is empty")

    with Image.open(image_paths[0]) as first_image:
        width_px, height_px = first_image.size
    if width_px <= 0 or height_px <= 0:
        raise ValueError("Invalid first image dimensions")

    prs = Presentation()
    blank_layout = prs.slide_layouts[6]

    slide_width = Emu(9144000)  # 10 inches
    slide_height = Emu(int(slide_width * height_px / width_px))
    prs.slide_width = slide_width
    prs.slide_height = slide_height

    if title:
        prs.core_properties.title = title

    for image_path in image_paths:
        slide = prs.slides.add_slide(blank_layout)
        slide.shapes.add_picture(image_path, 0, 0, width=prs.slide_width, height=prs.slide_height)

    prs.save(output_path)


def _normalize_text(value: object) -> str:
    if not isinstance(value, str):
        return ""
    return value.strip()


def _extract_slide_rows_for_master_meta(source_pptx: str) -> list[dict[str, object]]:
    try:
        with open(source_pptx, "rb") as fp:
            payload = fp.read()
    except Exception:
        return []
    try:
        context = extract_pptx_context(
            payload,
            filename=os.path.basename(source_pptx),
            source_url=None,
        )
    except Exception:
        return []

    raw_slides = context.get("slides") if isinstance(context, dict) else None
    if not isinstance(raw_slides, list):
        return []

    rows: list[dict[str, object]] = []
    for row in raw_slides:
        if not isinstance(row, dict):
            continue
        slide_number = row.get("slide_number")
        if not isinstance(slide_number, int) or slide_number <= 0:
            continue
        rows.append(
            {
                "slide_number": slide_number,
                "layout_name": _normalize_text(row.get("layout_name")) or None,
                "layout_placeholders": [
                    _normalize_text(item)
                    for item in (row.get("layout_placeholders") if isinstance(row.get("layout_placeholders"), list) else [])
                    if _normalize_text(item)
                ],
                "master_name": _normalize_text(row.get("master_name")) or None,
                "master_texts": [
                    _normalize_text(item)
                    for item in (row.get("master_texts") if isinstance(row.get("master_texts"), list) else [])
                    if _normalize_text(item)
                ],
            }
        )
    rows.sort(key=lambda row: int(row.get("slide_number") or 0))
    return rows


class _MasterDefinitionRow(BaseModel):
    slide_number: int = Field(...)
    layout_name: str | None = Field(default=None)
    layout_placeholders: list[str] = Field(default_factory=list)
    master_name: str | None = Field(default=None)
    master_texts: list[str] = Field(default_factory=list)


def _collect_unique_master_definition_rows(source_pptx: str) -> list[_MasterDefinitionRow]:
    rows = _extract_slide_rows_for_master_meta(source_pptx)
    if not rows:
        return []

    selected: list[_MasterDefinitionRow] = []
    seen_keys: set[tuple[str, str, tuple[str, ...]]] = set()
    for row in rows:
        master_name = _normalize_text(row.get("master_name")).lower()
        layout_name = _normalize_text(row.get("layout_name")).lower()
        placeholders = tuple(
            _normalize_text(item).lower()
            for item in (row.get("layout_placeholders") if isinstance(row.get("layout_placeholders"), list) else [])
            if _normalize_text(item)
        )
        if not master_name and not layout_name and not placeholders:
            slide_number = row.get("slide_number")
            if isinstance(slide_number, int):
                key = ("unknown", f"slide-{slide_number}", tuple())
            else:
                key = ("unknown", "unknown", tuple())
        else:
            key = (master_name, layout_name, placeholders)
        if key in seen_keys:
            continue
        seen_keys.add(key)
        slide_number_value = row.get("slide_number")
        if not isinstance(slide_number_value, int) or slide_number_value <= 0:
            continue
        selected.append(
            _MasterDefinitionRow(
                slide_number=slide_number_value,
                layout_name=_normalize_text(row.get("layout_name")) or None,
                layout_placeholders=[
                    _normalize_text(item)
                    for item in (row.get("layout_placeholders") if isinstance(row.get("layout_placeholders"), list) else [])
                    if _normalize_text(item)
                ],
                master_name=_normalize_text(row.get("master_name")) or None,
                master_texts=[
                    _normalize_text(item)
                    for item in (row.get("master_texts") if isinstance(row.get("master_texts"), list) else [])
                    if _normalize_text(item)
                ],
            )
        )
    return selected


def _fill_text_if_possible(shape: object, text: str) -> bool:
    if not hasattr(shape, "has_text_frame"):
        return False
    if not bool(getattr(shape, "has_text_frame")):
        return False
    text_frame = getattr(shape, "text_frame", None)
    if text_frame is None:
        return False
    try:
        text_frame.clear()
        text_frame.text = text
    except Exception:
        return False
    return True


def _clear_text_recursive(shape: object) -> None:
    _fill_text_if_possible(shape, "")
    child_shapes = getattr(shape, "shapes", None)
    if child_shapes is None:
        return
    try:
        for child in list(child_shapes):
            _clear_text_recursive(child)
    except Exception:
        return


def _clear_all_text_from_shape_collection(shape_collection: object) -> None:
    if shape_collection is None:
        return
    try:
        for shape in list(shape_collection):
            _clear_text_recursive(shape)
    except Exception:
        return


def _clear_master_definition_text(presentation: object) -> None:
    # スライドマスター画像には文字を残さない。
    slide_masters = getattr(presentation, "slide_masters", None)
    if slide_masters is None:
        return
    try:
        for master in list(slide_masters):
            _clear_all_text_from_shape_collection(getattr(master, "shapes", None))
            slide_layouts = getattr(master, "slide_layouts", None)
            if slide_layouts is None:
                continue
            for layout in list(slide_layouts):
                _clear_all_text_from_shape_collection(getattr(layout, "shapes", None))
    except Exception:
        return


def _render_pptx_to_images(
    *,
    source_pptx: str,
    out_dir: str,
    dpi_value: int,
    stem: str,
) -> list[str]:
    image_paths: list[str] = []
    with tempfile.TemporaryDirectory(prefix="pptx_render_", dir=os.path.dirname(out_dir)) as convert_dir:
        rendered_pdf = os.path.join(convert_dir, f"{stem}.pdf")
        convert_cmd = [
            "soffice",
            "--headless",
            "--convert-to",
            "pdf",
            "--outdir",
            convert_dir,
            source_pptx,
        ]
        convert_result = subprocess.run(convert_cmd, capture_output=True, text=True, check=False, timeout=120)
        if convert_result.returncode != 0:
            raise RuntimeError(
                "Failed to render PPTX with LibreOffice.\n"
                f"stdout: {convert_result.stdout}\n"
                f"stderr: {convert_result.stderr}"
            )
        if not os.path.isfile(rendered_pdf):
            raise RuntimeError(f"Converted PDF not found: {rendered_pdf}")

        pages = convert_from_path(rendered_pdf, dpi=dpi_value, fmt="png")
        for idx, page in enumerate(pages, start=1):
            image_path = os.path.join(out_dir, f"{stem}_master_{idx:02d}.png")
            page.save(image_path, format="PNG")
            image_paths.append(image_path)
    return image_paths


def _render_master_definition_images(
    *,
    source_pptx: str,
    out_dir: str,
    dpi_value: int,
    stem: str,
) -> tuple[list[str], list[dict[str, object]]]:
    try:
        from pptx import Presentation
    except Exception as exc:
        raise RuntimeError(f"python-pptx is required for master definition rendering: {exc}") from exc

    rows = _collect_unique_master_definition_rows(source_pptx)
    if not rows:
        return [], []

    presentation = Presentation(source_pptx)
    _clear_master_definition_text(presentation)
    layout_pool = list(presentation.slide_layouts)

    row_by_slide_number = {row.slide_number: row for row in rows}
    selected_layouts: list[tuple[_MasterDefinitionRow, object]] = []
    for slide_idx, slide in enumerate(presentation.slides, start=1):
        row = row_by_slide_number.get(slide_idx)
        if row is None:
            continue
        try:
            layout = slide.slide_layout
        except Exception:
            continue
        selected_layouts.append((row, layout))

    if not selected_layouts and layout_pool:
        fallback = rows[0]
        selected_layouts.append((fallback, layout_pool[0]))

    for _row, layout in selected_layouts:
        new_slide = presentation.slides.add_slide(layout)
        _clear_all_text_from_shape_collection(getattr(new_slide, "shapes", None))

    with tempfile.TemporaryDirectory(prefix="pptx_master_defs_", dir=os.path.dirname(out_dir)) as temp_dir:
        augmented_path = os.path.join(temp_dir, f"{stem}_master_defs.pptx")
        presentation.save(augmented_path)
        all_images = _render_pptx_to_images(
            source_pptx=augmented_path,
            out_dir=out_dir,
            dpi_value=dpi_value,
            stem=f"{stem}_master_defs",
        )

    definition_count = len(selected_layouts)
    if definition_count <= 0 or len(all_images) < definition_count:
        return [], []

    definition_images = all_images[-definition_count:]
    non_definition_images = all_images[:-definition_count]
    for path in non_definition_images:
        try:
            os.remove(path)
        except Exception:
            continue
    metadata: list[dict[str, object]] = []
    for image_path, (row, _layout) in zip(definition_images, selected_layouts):
        metadata.append(
            {
                "image_path": image_path,
                "source_title": None,
                "source_texts": [],
                "source_layout_name": row.layout_name,
                "source_layout_placeholders": list(row.layout_placeholders),
                "source_master_name": row.master_name,
                "source_master_texts": list(row.master_texts),
            }
        )
    return definition_images, metadata


@tool
def render_pptx_master_images_tool(
    pptx_path: Annotated[str, "Local path of PPTX file to render."],
    output_dir: Annotated[str, "Output directory for rendered images."] = "master_images",
    dpi: Annotated[int, "Rendering DPI (72-300)."] = 160,
    work_dir: Annotated[str, "Base working directory for relative paths."] = "/tmp",
    render_mode: Annotated[str, "Render mode: slide or master_definition."] = "slide",
):
    """
    Render PPTX into PNG files as local assets.
    This tool only handles local files. GCS I/O is handled by Data Analyst node.
    """
    try:
        safe_work_dir = os.path.abspath(work_dir or "/tmp")
        _ensure_dir(safe_work_dir)

        source_pptx = _resolve_path(pptx_path, safe_work_dir)
        if not os.path.isfile(source_pptx):
            return f"Error: PPTX file not found: {source_pptx}"

        out_dir = os.path.abspath(_resolve_path(output_dir, safe_work_dir))
        _ensure_dir(out_dir)

        dpi_value = max(72, min(300, int(dpi)))
        stem = Path(source_pptx).stem
        normalized_mode = (render_mode or "slide").strip().lower()
        image_paths: list[str] = []
        image_metadata: list[dict[str, object]] = []
        if normalized_mode == "master_definition":
            image_paths, image_metadata = _render_master_definition_images(
                source_pptx=source_pptx,
                out_dir=out_dir,
                dpi_value=dpi_value,
                stem=stem,
            )
        else:
            image_paths = _render_pptx_to_images(
                source_pptx=source_pptx,
                out_dir=out_dir,
                dpi_value=dpi_value,
                stem=stem,
            )

        return json.dumps(
            {
                "status": "ok",
                "pptx_path": source_pptx,
                "image_paths": image_paths,
                "image_metadata": image_metadata,
                "render_mode": normalized_mode,
            },
            ensure_ascii=False,
        )
    except Exception as e:
        logger.error("render_pptx_master_images_tool failed: %s", e)
        return f"Error: {e}"


@tool
def package_visual_assets_tool(
    image_paths: Annotated[list[str], "Local image paths to package in order."],
    output_basename: Annotated[str, "Base filename for pptx/pdf/zip outputs."] = "deck",
    output_dir: Annotated[str, "Output directory for packaged files."] = "packaged_assets",
    deck_title: Annotated[str, "Optional deck title embedded into PPTX metadata."] = "Generated Deck",
    work_dir: Annotated[str, "Base working directory for relative paths."] = "/tmp",
):
    """
    Package local images into PPTX, PDF, and ZIP files.
    This tool only handles local files. GCS I/O is handled by Data Analyst node.
    """
    try:
        safe_work_dir = os.path.abspath(work_dir or "/tmp")
        _ensure_dir(safe_work_dir)
        out_dir = os.path.abspath(_resolve_path(output_dir, safe_work_dir))
        _ensure_dir(out_dir)

        normalized_images: list[str] = []
        for raw_path in image_paths or []:
            if not isinstance(raw_path, str) or not raw_path.strip():
                continue
            resolved = _resolve_path(raw_path.strip(), safe_work_dir)
            if os.path.isfile(resolved):
                normalized_images.append(resolved)

        if not normalized_images:
            return "Error: No valid local images found for packaging."

        safe_basename = "".join(ch if ch.isalnum() or ch in {"-", "_"} else "_" for ch in output_basename).strip("_")
        safe_basename = safe_basename or "deck"

        pptx_path = os.path.join(out_dir, f"{safe_basename}.pptx")
        pdf_path = os.path.join(out_dir, f"{safe_basename}.pdf")
        zip_path = os.path.join(out_dir, f"{safe_basename}.zip")

        _build_pptx_from_images(normalized_images, pptx_path, title=deck_title)

        image_bytes_list: list[bytes] = []
        for image_path in normalized_images:
            with open(image_path, "rb") as fp:
                image_bytes_list.append(fp.read())
        pdf_bytes = assemble_pdf_from_images(image_bytes_list)
        with open(pdf_path, "wb") as fp:
            fp.write(pdf_bytes)

        with zipfile.ZipFile(zip_path, "w", compression=zipfile.ZIP_DEFLATED) as archive:
            archive.write(pptx_path, arcname=os.path.basename(pptx_path))
            archive.write(pdf_path, arcname=os.path.basename(pdf_path))
            for idx, image_path in enumerate(normalized_images, start=1):
                ext = Path(image_path).suffix or ".png"
                archive.write(image_path, arcname=f"slides/slide_{idx:02d}{ext}")

        return json.dumps(
            {
                "status": "ok",
                "input_image_count": len(normalized_images),
                "pptx_path": pptx_path,
                "pdf_path": pdf_path,
                "zip_path": zip_path,
            },
            ensure_ascii=False,
        )
    except Exception as e:
        logger.error("package_visual_assets_tool failed: %s", e)
        return f"Error: {e}"
