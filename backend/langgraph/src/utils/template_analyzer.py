"""
PPTXテンプレート分析の統合モジュール。

構造的抽出（python-pptx）とレイアウト別PNGレンダリングを組み合わせて、
統合的なDesignContextを生成する。
"""
import logging
import uuid
from typing import Dict, Optional

from src.schemas.design import DesignContext
from src.utils.pptx_extractor import extract_design_context_structure
from src.utils.pptx_renderer import PPTXRenderer, RenderedSlide
from src.utils.storage import upload_to_gcs

logger = logging.getLogger(__name__)


from io import BytesIO
from pptx import Presentation

def _prepare_visualization_pptx(pptx_bytes: bytes) -> bytes:
    """レンダリング用にPPTXを加工（全レイアウトの空スライドを生成）
    
    元のスライドを全て削除し、定義されている全てのレイアウトについて
    1枚ずつスライドを追加したPPTXを生成する。
    これにより、レンダリング結果のn番目の画像が確実にn番目のレイアウトに対応する。
    """
    prs = Presentation(BytesIO(pptx_bytes))
    
    # 既存のスライドを全て削除
    # xml_slides = prs.slides._sldIdLst  
    # slides = list(xml_slides)
    # for s in slides:
    #     xml_slides.remove(s)
    # 上記の内部API依存を避け、逆順で削除する安全な方法
    for i in range(len(prs.slides) - 1, -1, -1):
        rId = prs.slides._sldIdLst[i].rId
        prs.part.drop_rel(rId)
        del prs.slides._sldIdLst[i]
        
    # 全レイアウトについてスライドを追加
    # prs.slide_layouts はマスターのスライドレイアウト
    for layout in prs.slide_layouts:
        prs.slides.add_slide(layout)
        
    output = BytesIO()
    prs.save(output)
    return output.getvalue()


async def analyze_pptx_template(
    pptx_bytes: bytes,
    filename: str = "template.pptx",
    upload_to_gcs_enabled: bool = True,
    render_dpi: int = 150
) -> DesignContext:
    """
    PPTXテンプレートを分析し、レイアウト別のDesignContextを生成。
    
    処理フロー:
    1. 構造的抽出（python-pptx）: カラー、フォント、レイアウト情報
    2. レンダリング用PPTX生成: 各レイアウトのスライドを作成
    3. レイアウト別PNGレンダリング: 各レイアウトタイプごとにテンプレート画像を生成
    4. GCSアップロード（オプション）: 生成画像をCloudStorageにアップロード
    
    Args:
        pptx_bytes: PPTXファイルのバイトデータ
        filename: ファイル名（ログ用）
        upload_to_gcs_enabled: GCSにアップロードするかどうか
        render_dpi: レンダリング解像度（デフォルト: 150）
    
    Returns:
        DesignContext: レイアウト別画像を含む統合デザインコンテキスト
    """
    logger.info(f"Analyzing PPTX template: {filename}")
    
    # Step 1: 構造的抽出（元のPPTXから抽出）
    (color_scheme, font_scheme, layouts, background, 
     slide_master_count, layout_count) = extract_design_context_structure(pptx_bytes, filename)
    
    logger.info(
        f"Structural extraction complete: "
        f"{layout_count} layouts, "
        f"colors: {color_scheme.accent1}"
    )
    
    # 初期 DesignContext を作成（画像なし）
    design_context = DesignContext(
        color_scheme=color_scheme,
        font_scheme=font_scheme,
        layouts=layouts,
        background=background,
        source_filename=filename,
        slide_master_count=slide_master_count,
        layout_count=layout_count
    )
    
    # Step 2: レンダリング用PPTXの準備
    # 各レイアウトの見本スライドを含む一時的なPPTXを生成
    try:
        viz_pptx_bytes = _prepare_visualization_pptx(pptx_bytes)
        logger.info("Prepared visualization PPTX with layout slides")
    except Exception as e:
        logger.error(f"Failed to prepare visualization PPTX: {e}")
        # 失敗した場合は元のPPTXを使用（ただしマッピングがずれる可能性あり）
        viz_pptx_bytes = pptx_bytes
    
    # Step 3: 全スライド（レイアウト）をPNGにレンダリング
    try:
        renderer = PPTXRenderer(dpi=render_dpi)
        all_slides: list[RenderedSlide] = await renderer.render_to_png(viz_pptx_bytes)
        
        if all_slides:
            # レイアウトとレンダリング画像をマッピング
            # _prepare_visualization_pptx により、レンダリング順序 = レイアウト定義順序 が保証される
            import base64
            
            layout_images_base64: Dict[str, str] = {}
            layout_images: Dict[str, str] = {}
            
            for i, layout_info in enumerate(design_context.layouts):
                if i < len(all_slides):
                    rendered = all_slides[i]
                    layout_type = layout_info.layout_type
                    
                    # 同じタイプが複数ある場合は最初のものを使用
                    if layout_type not in layout_images_base64:
                        # Base64エンコードしてState保存用辞書に格納
                        b64_str = base64.b64encode(rendered.image_bytes).decode('utf-8')
                        layout_images_base64[layout_type] = b64_str
                        
                        logger.info(
                            f"Layout '{layout_type}' ({layout_info.name}) -> "
                            f"Slide {rendered.slide_number} ({rendered.width}x{rendered.height})"
                        )
                        
                        # GCSにアップロード (Optional)
                        if upload_to_gcs_enabled:
                            try:
                                # upload_to_gcs は自動でユニークなファイル名を生成する
                                url = upload_to_gcs(rendered.image_bytes, content_type="image/png")
                                layout_images[layout_type] = url
                                layout_info.template_image_url = url
                                logger.info(f"Uploaded template for {layout_type}: {url}")
                            except Exception as upload_err:
                                logger.warning(f"Failed to upload {layout_type}: {upload_err}")
            
            design_context.layout_images_base64 = layout_images_base64
            design_context.layout_images = layout_images
            
            # デフォルト画像は title_slide または最初のレイアウト
            if "title_slide" in layout_images_base64:
                design_context.default_template_image_base64 = layout_images_base64["title_slide"]
                design_context.default_template_image_url = layout_images.get("title_slide")
            elif layout_images_base64:
                first_type = list(layout_images_base64.keys())[0]
                design_context.default_template_image_base64 = layout_images_base64[first_type]
                design_context.default_template_image_url = layout_images.get(first_type)
            
            logger.info(f"Rendered {len(layout_images_base64)} unique layout types")
        else:
            logger.warning("No slides rendered from template")
            
    except Exception as e:
        logger.error(f"Template rendering failed: {e}")
        # レンダリング失敗しても構造的データは使用可能
    
    return design_context


async def analyze_pptx_template_simple(
    pptx_bytes: bytes,
    filename: str = "template.pptx"
) -> DesignContext:
    """構造的抽出のみを行うシンプル版（レンダリングなし）
    
    LibreOfficeがインストールされていない環境でも使用可能。
    
    Args:
        pptx_bytes: PPTXファイルのバイトデータ
        filename: ファイル名（ログ用）
        
    Returns:
        DesignContext: 構造的データのみのデザインコンテキスト
    """
    logger.info(f"Analyzing PPTX template (structure only): {filename}")
    
    (color_scheme, font_scheme, layouts, background, 
     slide_master_count, layout_count) = extract_design_context_structure(pptx_bytes, filename)
    
    return DesignContext(
        color_scheme=color_scheme,
        font_scheme=font_scheme,
        layouts=layouts,
        background=background,
        source_filename=filename,
        slide_master_count=slide_master_count,
        layout_count=layout_count
    )
